"""PPTX parser — extracts slide titles and visible text."""

from __future__ import annotations

import io
from typing import Optional

from agent.modules.knowledge.schemas import NormalizedDocument


def parse(
    raw: bytes,
    *,
    title: str = "",
    author: str = "",
    source_type: str = "project_doc",
    scope: str = "workspace",
    language: str = "zh",
    metadata: Optional[dict] = None,
) -> NormalizedDocument:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        return _empty(title, author, source_type, scope, language, metadata, f"pptx_parser_unavailable: {exc!r}")
    try:
        presentation = Presentation(io.BytesIO(raw))
    except Exception as exc:
        return _empty(title, author, source_type, scope, language, metadata, f"pptx_open_failed: {exc!r}")

    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text.append(shape.text.strip())
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("|", "\\|").replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        text.append("| " + " | ".join(cells) + " |")
        lines.append(f"## 幻灯片 {index}")
        lines.extend(text or ["（无可提取文字）"])
        lines.append("")
    return NormalizedDocument(
        title=title,
        author=author,
        source_type=source_type,
        scope=scope,
        language=language,
        format="pptx",
        normalized_markdown="\n".join(lines).strip(),
        metadata={**(metadata or {}), "format_hint": "pptx"},
        warnings=[],
    )


def _empty(title, author, source_type, scope, language, metadata, warning):
    return NormalizedDocument(
        title=title, author=author, source_type=source_type, scope=scope,
        language=language, format="pptx", normalized_markdown="",
        metadata=metadata or {}, warnings=[warning],
    )
